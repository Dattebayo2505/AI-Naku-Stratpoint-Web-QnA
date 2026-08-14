"""docparse.slides — LibreOffice conversion of .pptx decks into PDF.

Fixtures are *generated* zips rather than committed binaries, matching
tests/test_docparse_render.py: deterministic, tiny, and the construction
documents what "is a deck" actually means.

Nothing in this file requires LibreOffice to be installed. The converter is
injected; only the integration test in Task 6 spawns the real binary.
"""

from __future__ import annotations

import zipfile

import pytest

from stratpoint_rag.docparse import render, slides


# ── fixtures ────────────────────────────────────────────────────────────────


@pytest.fixture
def deck(tmp_path):
    """A minimal zip shaped like a .pptx: it carries ppt/presentation.xml."""
    path = tmp_path / "brief.pptx"
    with zipfile.ZipFile(path, "w") as z:
        z.writestr("[Content_Types].xml", "<Types/>")
        z.writestr("ppt/presentation.xml", "<presentation/>")
        z.writestr("ppt/slides/slide1.xml", "<sld/>")
    return path


@pytest.fixture
def docx(tmp_path):
    """A Word document renamed .pptx. Same magic bytes, not a deck."""
    path = tmp_path / "notadeck.pptx"
    with zipfile.ZipFile(path, "w") as z:
        z.writestr("[Content_Types].xml", "<Types/>")
        z.writestr("word/document.xml", "<document/>")
    return path


# ── detection ───────────────────────────────────────────────────────────────


def test_is_pptx_accepts_a_deck(deck):
    assert slides.is_pptx(deck) is True


def test_is_pptx_rejects_a_renamed_docx(docx):
    """Content decides, not the filename — the extension says otherwise here."""
    assert slides.is_pptx(docx) is False


def test_is_pptx_rejects_a_malformed_zip(tmp_path):
    path = tmp_path / "broken.pptx"
    path.write_bytes(b"PK\x03\x04" + b"\x00" * 64)
    assert slides.is_pptx(path) is False


def test_is_pptx_rejects_a_pdf(tmp_path):
    path = tmp_path / "brief.pdf"
    path.write_bytes(b"%PDF-1.7\n%%EOF\n")
    assert slides.is_pptx(path) is False


def test_is_pptx_rejects_a_missing_file(tmp_path):
    """A path that does not exist is not a deck; it must not raise."""
    assert slides.is_pptx(tmp_path / "gone.pptx") is False


def test_conversion_failed_is_an_unsupported_document():
    """Subclassing is load-bearing: app.py's existing 400 handler catches
    UnsupportedDocument, so ConversionFailed needs no new wiring there."""
    assert issubclass(slides.ConversionFailed, render.UnsupportedDocument)


# ── conversion ──────────────────────────────────────────────────────────────


def _fake_converter(calls):
    """A Converter that records its calls and writes a plausible PDF."""

    def convert(src, outdir):
        calls.append((src, outdir))
        (outdir / f"{src.stem}.pdf").write_bytes(b"%PDF-1.7\n%%EOF\n")

    return convert


def test_ensure_pdf_converts_a_deck(deck):
    calls = []
    out = slides.ensure_pdf(deck, convert=_fake_converter(calls))

    assert out == deck.parent / slides.CONVERTED_NAME
    assert out.is_file()
    assert len(calls) == 1


def test_ensure_pdf_is_a_no_op_for_a_pdf(tmp_path):
    """Both call sites invoke it unconditionally, so a non-deck must pass
    straight through without spawning anything."""
    path = tmp_path / "brief.pdf"
    path.write_bytes(b"%PDF-1.7\n%%EOF\n")
    calls = []

    assert slides.ensure_pdf(path, convert=_fake_converter(calls)) == path
    assert calls == []


def test_ensure_pdf_caches_the_conversion(deck):
    """Called at upload and again at parse. The second call is a stat, not a
    ten-second subprocess."""
    calls = []
    convert = _fake_converter(calls)

    first = slides.ensure_pdf(deck, convert=convert)
    second = slides.ensure_pdf(deck, convert=convert)

    assert first == second
    assert len(calls) == 1


def test_ensure_pdf_reconverts_when_the_cache_is_empty(deck):
    """A zero-byte converted.pdf is a crashed previous run, not a cache hit."""
    (deck.parent / slides.CONVERTED_NAME).write_bytes(b"")
    calls = []

    slides.ensure_pdf(deck, convert=_fake_converter(calls))

    assert len(calls) == 1


def test_ensure_pdf_stages_the_deck_outside_the_upload_directory(deck):
    """LibreOffice writes beside its input (a .~lock.<name># at minimum), and
    on Windows the upload directory routinely sits inside a folder Controlled
    Folder Access protects, where Defender blocks soffice.bin from writing at
    all. The converter must therefore never see the stored upload's path."""
    calls = []
    record = _fake_converter(calls)
    seen = {}

    def convert(src, outdir):
        # Read while the temp tree still exists — ensure_pdf tears it down.
        seen["bytes"] = src.read_bytes()
        record(src, outdir)

    out = slides.ensure_pdf(deck, convert=convert)

    src, outdir = calls[0]
    assert src != deck
    assert src.name == deck.name  # stem drives the produced filename
    assert deck.parent not in src.parents
    assert seen["bytes"] == deck.read_bytes()
    assert outdir != src.parent  # nothing dropped beside the input is output
    assert out.read_bytes().startswith(b"%PDF")


def test_ensure_pdf_cleans_up_the_staged_copy(deck):
    """The temp tree goes away with the context manager; only converted.pdf
    is left behind, inside the upload's own directory."""
    calls = []

    slides.ensure_pdf(deck, convert=_fake_converter(calls))

    staged, _ = calls[0]
    assert not staged.exists()
    assert sorted(p.name for p in deck.parent.iterdir()) == sorted(
        [deck.name, slides.CONVERTED_NAME]
    )


def test_ensure_pdf_raises_when_nothing_was_produced(deck):
    """soffice returns 0 on inputs it silently declined, so the exit code is
    never the check — the output file is."""

    def convert(src, outdir):
        return None  # exits cleanly, writes nothing

    with pytest.raises(slides.ConversionFailed) as ex:
        slides.ensure_pdf(deck, convert=convert)
    assert deck.name in str(ex.value)


def test_ensure_pdf_raises_when_the_output_is_empty(deck):
    def convert(src, outdir):
        (outdir / f"{src.stem}.pdf").write_bytes(b"")

    with pytest.raises(slides.ConversionFailed):
        slides.ensure_pdf(deck, convert=convert)


def test_ensure_pdf_wraps_a_converter_crash(deck):
    """A subprocess failure must arrive as ConversionFailed, which app.py
    already maps to 400 — never as a raw OSError escaping to a 500."""

    def convert(src, outdir):
        raise OSError("soffice died")

    with pytest.raises(slides.ConversionFailed):
        slides.ensure_pdf(deck, convert=convert)


# ── the soffice command line ────────────────────────────────────────────────


def test_soffice_command_isolates_the_user_profile(monkeypatch, deck, tmp_path):
    """The regression guard for the silent failure.

    Without -env:UserInstallation, an already-running soffice — or a second
    concurrent upload — makes this invocation exit 0 having converted nothing.
    It looks exactly like success, which is why it is asserted here rather than
    trusted to a comment.
    """
    monkeypatch.delenv("SOFFICE_TIMEOUT", raising=False)
    seen = {}

    def fake_run(cmd, **kwargs):
        seen["cmd"] = cmd
        seen["kwargs"] = kwargs

        class R:
            returncode = 0

        return R()

    monkeypatch.setattr(slides, "find_soffice", lambda: "soffice")
    monkeypatch.setattr(slides.subprocess, "run", fake_run)

    slides._soffice_convert(deck, tmp_path)

    cmd = seen["cmd"]
    assert any(a.startswith("-env:UserInstallation=file://") for a in cmd)
    assert "--headless" in cmd
    assert "--convert-to" in cmd
    assert cmd[cmd.index("--convert-to") + 1] == "pdf"
    assert cmd[cmd.index("--outdir") + 1] == str(tmp_path)
    assert cmd[-1] == str(deck)
    assert seen["kwargs"]["timeout"] == 120


def test_soffice_convert_raises_on_a_timeout(monkeypatch, deck, tmp_path):
    """Under a hung conversion the upload must fail cleanly, not hang."""
    monkeypatch.delenv("SOFFICE_TIMEOUT", raising=False)

    def fake_run(cmd, **kwargs):
        raise slides.subprocess.TimeoutExpired(cmd, 120)

    monkeypatch.setattr(slides, "find_soffice", lambda: "soffice")
    monkeypatch.setattr(slides.subprocess, "run", fake_run)

    with pytest.raises(slides.ConversionFailed) as ex:
        slides._soffice_convert(deck, tmp_path)
    assert "120" in str(ex.value)


def test_find_soffice_prefers_the_configured_binary(monkeypatch, tmp_path):
    binary = tmp_path / "soffice"
    binary.write_text("#!/bin/sh\n")
    monkeypatch.setenv("SOFFICE_BINARY", str(binary))

    assert slides.find_soffice() == str(binary)


def test_find_soffice_raises_when_absent(monkeypatch):
    """RuntimeError, not ConversionFailed: the file is fine, the server is
    misconfigured. api/app.py maps that to 503."""
    monkeypatch.setenv("SOFFICE_BINARY", "")
    monkeypatch.setattr(slides.shutil, "which", lambda name: None)
    monkeypatch.setattr(slides, "_WINDOWS_CANDIDATES", ())

    with pytest.raises(RuntimeError) as ex:
        slides.find_soffice()
    assert "LibreOffice" in str(ex.value)


# ── the pipeline entry point ────────────────────────────────────────────────


def test_open_brief_marks_a_converted_deck_as_slides(deck, monkeypatch):
    """The whole feature turns on this flag: without it the converted PDF's
    text layer sends every slide down the free text route."""
    import pymupdf

    def convert(src, outdir):
        doc = pymupdf.open()
        doc.new_page(width=720, height=405)
        doc.save(outdir / f"{src.stem}.pdf")
        doc.close()

    with slides.open_brief(deck, convert=convert) as doc:
        assert doc.kind == "slides"
        assert doc.page_count == 1


def test_open_brief_leaves_a_pdf_alone(tmp_path):
    import pymupdf

    path = tmp_path / "brief.pdf"
    doc = pymupdf.open()
    doc.new_page(width=595, height=842)
    doc.save(path)
    doc.close()

    with slides.open_brief(path) as opened:
        assert opened.kind == "pdf"


# ── live conversion ─────────────────────────────────────────────────────────


@pytest.mark.integration
def test_real_libreoffice_converts_a_real_deck(tmp_path):
    """The only test here that spawns LibreOffice. Deselected by default via
    pyproject's addopts, like the crawler's live test.

    Guards what the fake converter structurally cannot: that the command line
    is one LibreOffice actually accepts, and that a real deck yields one PDF
    page per slide.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for n in (1, 2, 3):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Slide {n}"
        slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(6), Inches(1)
        ).text_frame.text = "Migrate to Kubernetes on AWS."
    path = tmp_path / "live.pptx"
    prs.save(path)

    with slides.open_brief(path) as doc:
        assert doc.kind == "slides"
        assert doc.page_count == 3
        image = doc.rasterize(0)

    assert image.startswith(b"\xff\xd8\xff")  # JPEG
    assert len(image) > 1000
